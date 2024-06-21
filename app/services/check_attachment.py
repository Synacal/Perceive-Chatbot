from fastapi import FastAPI, HTTPException
import base64
from io import BytesIO
from PyPDF2 import PdfReader
from app.utils.prompts import questions, prompts
from app.core.azure_client import client
import json
from app.utils.prompts import questions, prompts
from pptx import Presentation
from docx import Document
import requests
from bs4 import BeautifulSoup
from requests.exceptions import RequestException

# from pptx import Presentation
# from docx import Document


def get_pdf_content(attachment_base64: str) -> str:
    try:
        # Decode the base64 string
        pdf_bytes = base64.b64decode(attachment_base64)
        # Use BytesIO to read the PDF bytes
        pdf_file = BytesIO(pdf_bytes)
        # Create a PDF reader
        pdf_reader = PdfReader(pdf_file)
        # Extract text from each page
        content = ""
        for page in pdf_reader.pages:
            content += page.extract_text() + "\n"

        return content
    except Exception as e:
        raise HTTPException(
            status_code=400, detail=f"Error reading PDF content: {str(e)}"
        )


def get_pptx_content(attachment_base64: str) -> str:
    try:
        # Decode the base64 string
        pptx_bytes = base64.b64decode(attachment_base64)
        # Use BytesIO to read the pptx bytes
        pptx_file = BytesIO(pptx_bytes)
        # Open the presentation
        presentation = Presentation(pptx_file)
        # Extract text from each slide
        content = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
        return content
    except Exception as e:
        raise HTTPException(
            status_code=400, detail=f"Error reading PPTX content: {str(e)}"
        )


def get_docx_content(attachment_base64: str) -> str:
    try:
        # Decode the base64 string
        docx_bytes = base64.b64decode(attachment_base64)
        # Use BytesIO to read the docx bytes
        docx_file = BytesIO(docx_bytes)
        # Open the document
        document = Document(docx_file)
        # Extract text from each paragraph
        content = ""
        for paragraph in document.paragraphs:
            content += paragraph.text + "\n"
        return content
    except Exception as e:
        raise HTTPException(
            status_code=400, detail=f"Error reading DOCX content: {str(e)}"
        )


async def get_txt_content(attachment_base64: str) -> str:
    try:
        # Decode the base64 string
        txt_bytes = base64.b64decode(attachment_base64)
        # Convert bytes to string (assuming UTF-8 encoding)
        content = txt_bytes.decode("utf-8")
        return content
    except Exception as e:
        raise HTTPException(
            status_code=400, detail=f"Error reading TXT content: {str(e)}"
        )


async def get_content(attachments: list) -> str:
    try:
        # Get the content from the attachments
        content = ""
        for attachment in attachments:
            if attachment.fileType == "application/pdf":
                content += get_pdf_content(attachment.file)
            elif (
                attachment.fileType
                == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ):
                content += get_pptx_content(attachment.file)
            elif (
                attachment.fileType == "application/msword"
                or attachment.fileType
                == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ):
                content += get_docx_content(attachment.file)
            elif attachment.fileType == "text/plain":
                content += get_txt_content(attachment.file)
            else:
                raise HTTPException(
                    status_code=400,
                    detail="Attachment type not supported. Supported types are pdf, pptx,txt, and docx.",
                )
        return content
    except Exception as e:
        raise HTTPException(
            status_code=400, detail=f"Error getting content from attachments: {str(e)}"
        )


async def get_web_content(web_urls: list) -> str:
    try:
        # Get the content from the web URLs
        content = ""
        for url in web_urls:
            try:
                response = requests.get(url)
                response.raise_for_status()  # Raise an exception for HTTP errors

                soup = BeautifulSoup(response.text, "html.parser")
                text = soup.get_text()  # Get the text content from the HTML

                content += text + "\n\n"  # Add the extracted text to the content
            except RequestException as req_ex:
                print(f"Error fetching URL '{url}': {str(req_ex)}")
                # Log the error or handle it as needed

        return content
    except Exception as e:
        raise HTTPException(
            status_code=400, detail=f"Error getting content from web URLs: {str(e)}"
        )


def get_questions(category_id: str) -> list:
    try:
        # Get the questions from the category
        if category_id == "1" or category_id == "2":
            selectedQuestions = questions[0:2] + questions[5:13]
        elif category_id == "3":
            selectedQuestions = questions[13:26]
        elif category_id == "4":
            selectedQuestions = questions[26:35]
        elif category_id == "5":
            selectedQuestions = questions[35:42]
        elif category_id == "6":
            selectedQuestions = questions[0:4] + questions[42:55]
        elif category_id == "7":
            selectedQuestions = questions[0:4] + questions[55:61]
        elif category_id == "8":
            selectedQuestions = questions[0:4] + questions[61:72]
        elif category_id == "9":
            selectedQuestions = questions[0:4] + questions[72:80]
        elif category_id == "10":
            selectedQuestions = questions[0:4] + questions[80:89]
        else:
            selectedQuestions = questions[0:4]
        return selectedQuestions
    except Exception as e:
        raise HTTPException(
            status_code=400, detail=f"Error getting questions: {str(e)}"
        )


def get_prompts(category_id: str) -> list:
    try:
        # Get the prompts from the category
        if category_id == "1" or category_id == "2":
            selectedPrompts = prompts[0:2] + prompts[5:13]
        elif category_id == "3":
            selectedPrompts = prompts[13:26]
        elif category_id == "4":
            selectedPrompts = prompts[26:35]
        elif category_id == "5":
            selectedPrompts = prompts[35:42]
        elif category_id == "6":
            selectedPrompts = prompts[0:4] + prompts[42:55]
        elif category_id == "7":
            selectedPrompts = prompts[0:4] + prompts[55:61]
        elif category_id == "8":
            selectedPrompts = prompts[0:4] + prompts[61:72]
        elif category_id == "9":
            selectedPrompts = prompts[0:4] + prompts[72:80]
        elif category_id == "10":
            selectedPrompts = prompts[0:4] + prompts[80:89]
        else:
            selectedPrompts = prompts[0:4]
        return selectedPrompts
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error getting prompts: {str(e)}")


async def check_user_attachment_temp(
    questions: list,
    prompts: list,
    content: str,
    report_id: str,
    user_id: str,
    category_id: str,
):
    print(json.dumps(questions))
    print(json.dumps(prompts))

    system_prompt = f"""
        Given the user's response to the questions: {json.dumps(questions)},

        evaluate the completeness based on these criteria and provide the response always in JSON format as given below:
        {json.dumps(prompts)}

        Questions and completeness criteria are entered in order to match each other.

        The response should encapsulate all specified points to be considered complete.

        If the user's input lacks any required details, is ambiguous, or misses critical information, the output should be:
        "responses": [
            {{ "index list  of the question that has not been fully answered" }},
            ...
        ]

        example response (if question 1 and 3 are not answered fully):
        {{
            "responses": [
                1,
                3
            ]
        }}

        questions that have not been fully answered should be included in the responses key of the output JSON format.
        """

    message_text = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": content},
    ]

    try:
        completion = client.chat.completions.create(
            model="gpt-35-turbo",
            messages=message_text,
            temperature=0.7,
            max_tokens=800,
            top_p=0.95,
            frequency_penalty=0,
            presence_penalty=0,
            stop=None,
        )

        content = completion.choices[0].message.content
        print(content)

        if content:
            try:
                generated_content_json = content
                return generated_content_json
            except json.JSONDecodeError:
                return {
                    "status": "Error of code",
                }
        else:
            return {"Error of code"}
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Error generating response: {str(e)}"
        )


def find_question_number(question: str) -> int:
    try:
        question_number = questions.index(question) + 1
        return question_number
    except Exception as e:
        raise HTTPException(
            status_code=400, detail=f"Error finding question number: {str(e)}"
        )


async def check_user_attachment(
    answeredQuestion: str,
    checkPrompt: str,
    content: str,
    requirement_gathering_id,
    user_id: str,
    report_id: str,
):

    system_prompt = f"""Given the user's response to the question: '{answeredQuestion}',
            evaluate the completeness based on these criteria and provide the response always in JSON format as given below:
            '{checkPrompt}'

            The response should encapsulate all specified points to be considered complete.

            If the user's input lacks any required details, is ambiguous, or misses critical information, the output should be:
            {{"status": "false", 
              "question": "<appropriate follow-up question from the predefined list>"}}
            
            This indicates the need for additional information to fulfill the request comprehensively.

            For every gap identified in the user's response, select a follow-up question that precisely targets 
            the missing information. This could involve requesting more elaborate explanations, clarification of 
            technical terms, specific instances of technology application, or arguments supporting the novelty and 
            uniqueness of the concept. It should be included in the question key of the output JSON format. 
            Always format the output in JSON, including 'status' and 'question' keys, to streamline the evaluation 
            process and guide the user towards providing a fully rounded response.

            Conversely, if the user's answer meets all the outlined criteria, confirm the completeness with:
            {{"status": "true", 
            "question": "",
            "answer": "<the user's response to the question>"}}
            
            Avoid generating apology messages or phrases like "I'm sorry" in the follow-up questions or responses.

            Always format the output in JSON, including 'status' and 'question' keys, to streamline the evaluation 
            process and guide the user towards providing a fully rounded response.
            """

    message_text = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": content},
    ]

    try:
        completion = client.chat.completions.create(
            model="gpt-35-turbo",
            messages=message_text,
            temperature=0.7,
            max_tokens=800,
            top_p=0.95,
            frequency_penalty=0,
            presence_penalty=0,
            stop=None,
        )

        content = completion.choices[0].message.content
        print(content)

        if content:
            try:
                generated_content_json = json.loads(content)
                if (
                    "status" not in generated_content_json
                    or "question" not in generated_content_json
                ):
                    return {
                        "status": "false",
                        "question": f"I couldn't identify the required details in your response to the question: '{answeredQuestion}'. Can you provide more specific information or elaborate further?",
                    }
                return generated_content_json
            except json.JSONDecodeError:
                generated_content_json = {
                    "status": "false",
                    "question": f"I couldn't identify the required details in your response to the question: '{answeredQuestion}'. Can you provide more specific information or elaborate further?",
                }
        else:
            generated_content_json = {
                "status": "false",
                "question": f"I couldn't identify the required details in your response to the question: '{answeredQuestion}'. Can you provide more specific information or elaborate further?",
            }

        return generated_content_json

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
